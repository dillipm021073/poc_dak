#!/usr/bin/env python3
"""Generate D.A.K POC Proposal PPTX — matching the app mockup UI/UX"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# ── Design tokens (from the mockup) ──
BG = RGBColor(0xE8, 0xEC, 0xF4)          # soft lavender-gray background
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CARD_BORDER = RGBColor(0xE2, 0xE8, 0xF0)  # subtle card border
ACCENT = RGBColor(0x25, 0x63, 0xEB)       # primary blue
ACCENT_LIGHT = RGBColor(0xDB, 0xEA, 0xFE) # light blue tint
DARK = RGBColor(0x1F, 0x29, 0x37)         # headings
TEXT = RGBColor(0x37, 0x41, 0x51)          # body text
SUBTLE = RGBColor(0x6B, 0x72, 0x80)       # secondary text
MUTED = RGBColor(0x9C, 0xA3, 0xAF)        # muted text
GREEN = RGBColor(0x10, 0xB9, 0x81)
GREEN_LIGHT = RGBColor(0xD1, 0xFA, 0xE5)
PURPLE_LIGHT = RGBColor(0xE9, 0xD5, 0xFF)
AMBER_LIGHT = RGBColor(0xFE, 0xF3, 0xC7)
PINK_LIGHT = RGBColor(0xFC, 0xE7, 0xF3)
FONT = 'Segoe UI'

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def slide_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def card(slide, l, t, w, h, fill=WHITE, border=CARD_BORDER, radius=True):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        Inches(l), Inches(t), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if border:
        shp.line.color.rgb = border
        shp.line.width = Pt(0.75)
    else:
        shp.line.fill.background()
    return shp


def rect(slide, l, t, w, h, fill, border=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if border:
        shp.line.color.rgb = border
        shp.line.width = Pt(0.5)
    else:
        shp.line.fill.background()
    return shp


def txt(slide, l, t, w, h, text, size=14, color=DARK, bold=False,
        align=PP_ALIGN.LEFT, font=FONT):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font
    p.alignment = align
    return tb


def rich_line(slide, l, t, w, h, parts):
    """parts = [(text, size, color, bold), ...]"""
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    for text, size, color, bold in parts:
        r = p.add_run()
        r.text = text
        r.font.size = Pt(size)
        r.font.color.rgb = color
        r.font.bold = bold
        r.font.name = FONT
    return tb


def nav_bar(slide):
    """White top nav bar like the mockup"""
    card(slide, 0.4, 0.25, 12.533, 0.6, WHITE, CARD_BORDER, radius=True)
    # Star of David + D.A.K
    txt(slide, 0.7, 0.32, 0.3, 0.4, "✡", 18, ACCENT)
    txt(slide, 1.0, 0.3, 1.2, 0.5, "D.A.K", 20, DARK, bold=True)
    # Right side subtitle
    txt(slide, 9.5, 0.38, 3.3, 0.35, "POC Development Proposal", 11, MUTED, align=PP_ALIGN.RIGHT)


def footer(slide, page_num):
    txt(slide, 0.8, 7.1, 3, 0.3, "D.A.K — Digital Access Key", 9, MUTED)
    txt(slide, 10, 7.1, 2.8, 0.3, f"Slide {page_num} of 3", 9, MUTED, align=PP_ALIGN.RIGHT)


def check_circle(slide, l, t, size=0.18):
    c = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(l), Inches(t), Inches(size), Inches(size))
    c.fill.solid()
    c.fill.fore_color.rgb = GREEN
    c.line.fill.background()
    p = c.text_frame.paragraphs[0]
    p.text = "✓"
    p.font.size = Pt(7)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER


def section_header(slide, l, t, title, subtitle=None):
    """Blue underlined tab-style header like the mockup"""
    txt(slide, l, t, 5, 0.4, title, 26, DARK, bold=True)
    # blue underline
    rect(slide, l, t + 0.42, 1.8, 0.035, ACCENT)
    if subtitle:
        txt(slide, l, t + 0.5, 11, 0.35, subtitle, 13, SUBTLE)


# =====================================================================
# SLIDE 1 — SCOPE OF POC
# =====================================================================
s1 = prs.slides.add_slide(prs.slide_layouts[6])
slide_bg(s1)
nav_bar(s1)

section_header(s1, 0.8, 1.05, "Scope of POC",
               "Multi-portal faith-community platform with Spiritual AI Guru Bot, VPS deployment & Docker orchestration")

# ── Left: Platform Architecture (white card) ──
card(s1, 0.5, 1.75, 6.0, 5.2)
txt(s1, 0.8, 1.9, 5, 0.35, "Platform Architecture", 16, DARK, bold=True)
rect(s1, 0.8, 2.2, 0.7, 0.025, ACCENT)

arch_items = [
    ("Admin Portal", "Platform admin — approve communities, view GMV, manage users, payments"),
    ("Institute Portal", "Community admin — branding, events, live streams, messaging with members"),
    ("User Portal", "End user — community overview, live stream, support widget, events, calendar"),
    ("API Backend", "Node.js/Express + PostgreSQL — JWT auth, role-based access, 16+ endpoints"),
    ("Spiritual AI Guru Bot", "AI-powered spiritual chatbot — faith-specific guidance, scripture references"),
    ("VPS & Docker Deploy", "Production VPS, Docker orchestration, SSL/TLS, domain, monitoring"),
]

y = 2.4
for title, desc in arch_items:
    # light blue accent strip
    rect(s1, 0.8, y, 0.05, 0.52, ACCENT)
    card(s1, 0.85, y, 5.45, 0.52, RGBColor(0xF8, 0xFA, 0xFC), CARD_BORDER)
    rich_line(s1, 1.0, y + 0.03, 5.1, 0.22,
              [(title, 12, DARK, True)])
    txt(s1, 1.0, y + 0.26, 5.1, 0.22, desc, 9, SUBTLE)
    y += 0.6

# ── Right: Key Features (white card) ──
card(s1, 6.8, 1.75, 6.1, 5.2)
txt(s1, 7.1, 1.9, 5, 0.35, "Key Features in Scope", 16, DARK, bold=True)
rect(s1, 7.1, 2.2, 0.7, 0.025, ACCENT)

features = [
    ("Controlled Access", "QR/Invite-only onboarding, Chinese-wall isolation, 4 community types"),
    ("Live Streaming", "Embedded player, LIVE badge, viewer count, playback controls, recording gating"),
    ("Support & Payments", "One-time/Monthly, $16/$36/$100/$250 presets, 5% platform fee"),
    ("Events & Calendar", "Event management, upcoming events sidebar, mini calendar, ICS export"),
    ("Spiritual AI Guru Bot", "Faith-specific AI chatbot, scripture Q&A, community-context aware"),
    ("Private Messaging", "Threaded User↔Admin inbox, real-time, block capability"),
    ("Community Management", "Onboarding, branding, member management, stream scheduling"),
    ("Admin Dashboard", "GMV, users, subscribers metrics, community approve/suspend"),
    ("VPS Production Setup", "Server provisioning, Docker deploy, SSL, domain, backups"),
    ("Responsive UI", "Apple-inspired design, two-column layout, Star of David branding"),
]

y = 2.4
for title, desc in features:
    check_circle(s1, 7.1, y + 0.06)
    rich_line(s1, 7.35, y + 0.02, 5.3, 0.2,
              [(title, 11, DARK, True), ("  —  " + desc, 9, SUBTLE, False)])
    y += 0.4

footer(s1, 1)


# =====================================================================
# SLIDE 2 — TIMELINE (8-10 WEEKS)
# =====================================================================
s2 = prs.slides.add_slide(prs.slide_layouts[6])
slide_bg(s2)
nav_bar(s2)

section_header(s2, 0.8, 1.05, "Timeline",
               "8–10 week development plan including VPS setup & Spiritual AI Guru Bot — single full-stack developer")

# ── 5 phase cards ──
phases = [
    {
        "label": "WEEK 1–2",
        "title": "VPS Setup &\nInfrastructure",
        "tint": ACCENT_LIGHT,
        "items": [
            "VPS provisioning & config",
            "Docker environment setup",
            "PostgreSQL schema & seeds",
            "SSL/TLS & domain setup",
            "CI/CD pipeline basics",
            "JWT auth & OAuth",
            "Role-based access control",
        ]
    },
    {
        "label": "WEEK 3–4",
        "title": "Core Platform\nDevelopment",
        "tint": GREEN_LIGHT,
        "items": [
            "3 portal UIs (Admin,\n  Institute, User)",
            "Community management",
            "Events & calendar system",
            "Live streaming integration",
            "Messaging system",
            "Notifications engine",
        ]
    },
    {
        "label": "WEEK 5–6",
        "title": "Payments &\nAccess System",
        "tint": AMBER_LIGHT,
        "items": [
            "Stripe sandbox integration",
            "One-time & recurring flows",
            "Support widget with presets",
            "Platform fee (sliding scale)",
            "Access duration logic",
            "Webhook & refund handling",
        ]
    },
    {
        "label": "WEEK 7–8",
        "title": "Spiritual AI\nGuru Bot",
        "tint": PURPLE_LIGHT,
        "items": [
            "AI model selection & setup",
            "Faith-specific training data",
            "Context-aware responses",
            "Scripture reference engine",
            "Chat UI integration",
            "Content safety guardrails",
        ]
    },
    {
        "label": "WEEK 9–10",
        "title": "Testing, Polish\n& Demo",
        "tint": PINK_LIGHT,
        "items": [
            "End-to-end testing",
            "Security audit & hardening",
            "UI/UX polish to mockup",
            "Production deployment",
            "Documentation & handover",
            "Stakeholder live demo",
        ]
    },
]

cw = 2.3
cg = 0.12
cx = 0.5
for i, ph in enumerate(phases):
    x = cx + i * (cw + cg)
    # Main card
    card(s2, x, 1.85, cw, 4.35)

    # Tinted header area
    rect(s2, x, 1.85, cw, 0.9, ph["tint"], CARD_BORDER)

    # Week badge
    badge = card(s2, x + 0.12, 1.97, 0.95, 0.28, ACCENT, border=None)
    p = badge.text_frame.paragraphs[0]
    p.text = ph["label"]
    p.font.size = Pt(9)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = FONT
    p.alignment = PP_ALIGN.CENTER

    # Phase title
    txt(s2, x + 0.12, 2.32, cw - 0.24, 0.45, ph["title"], 13, DARK, bold=True)

    # Items
    iy = 2.85
    for item in ph["items"]:
        # small blue dot
        dot = s2.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 0.15), Inches(iy + 0.06), Inches(0.06), Inches(0.06))
        dot.fill.solid()
        dot.fill.fore_color.rgb = ACCENT
        dot.line.fill.background()
        txt(s2, x + 0.28, iy, cw - 0.4, 0.22, item, 10, TEXT)
        iy += 0.28

    # Arrow between cards
    if i < 4:
        ax = x + cw + 0.01
        arr = s2.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(ax), Inches(3.8), Inches(0.1), Inches(0.2))
        arr.fill.solid()
        arr.fill.fore_color.rgb = ACCENT
        arr.line.fill.background()

# ── Milestones card ──
card(s2, 0.5, 6.35, 12.4, 0.6)
txt(s2, 0.8, 6.4, 2, 0.25, "KEY MILESTONES", 10, ACCENT, bold=True)

milestones = [
    ("Week 2:", " VPS live, auth working"),
    ("Week 4:", " All 3 portals functional"),
    ("Week 6:", " Payments end-to-end"),
    ("Week 8:", " AI Guru Bot integrated"),
    ("Week 10:", " Live demo to stakeholders"),
]
mx = 0.8
for label, desc in milestones:
    rich_line(s2, mx, 6.62, 2.3, 0.25,
              [(label, 10, DARK, True), (desc, 10, SUBTLE, False)])
    mx += 2.45

footer(s2, 2)


# =====================================================================
# SLIDE 3 — COST BREAKDOWN (~$10K)
# =====================================================================
s3 = prs.slides.add_slide(prs.slide_layouts[6])
slide_bg(s3)
nav_bar(s3)

section_header(s3, 0.8, 1.05, "Cost Breakdown",
               "Development, infrastructure & operational costs for the 8–10 week POC phase")

def cost_card(slide, x, y, w, h, title, rows, total_lbl, total_val, highlight=False):
    bg = RGBColor(0xEF, 0xF6, 0xFF) if highlight else WHITE
    bdr = ACCENT if highlight else CARD_BORDER
    card(slide, x, y, w, h, bg, bdr)

    # Title row with tinted header
    tint = RGBColor(0xDB, 0xEA, 0xFE) if not highlight else RGBColor(0xCF, 0xDF, 0xFD)
    rect(slide, x, y, w, 0.45, tint, bdr)
    txt(slide, x + 0.2, y + 0.08, w - 0.4, 0.3, title, 11, ACCENT, bold=True)

    ry = y + 0.55
    for label, value in rows:
        txt(slide, x + 0.2, ry, w * 0.58, 0.25, label, 11, TEXT)
        txt(slide, x + w * 0.52, ry, w * 0.42, 0.25, value, 11, DARK, bold=True, align=PP_ALIGN.RIGHT)
        # separator
        rect(slide, x + 0.2, ry + 0.27, w - 0.4, 0.005, CARD_BORDER)
        ry += 0.32

    # Total
    ry += 0.05
    rect(slide, x + 0.2, ry, w - 0.4, 0.02, ACCENT)
    ry += 0.08
    txt(slide, x + 0.2, ry, w * 0.5, 0.3, total_lbl, 14, ACCENT, bold=True)
    txt(slide, x + w * 0.45, ry, w * 0.5, 0.3, total_val, 16, ACCENT, bold=True, align=PP_ALIGN.RIGHT)


# Card 1: Development
cost_card(s3, 0.5, 1.75, 3.9, 3.05,
    "DEVELOPMENT COST (ONE-TIME)",
    [
        ("Wk 1-2: VPS Setup & Infra", "$1,500"),
        ("Wk 3-4: Core Platform Dev", "$2,500"),
        ("Wk 5-6: Payments & Access", "$1,500"),
        ("Wk 7-8: Spiritual AI Guru Bot", "$2,000"),
        ("Wk 9-10: Testing & Demo", "$1,000"),
    ],
    "Development Total", "$8,500"
)

# Card 2: Infrastructure
cost_card(s3, 4.7, 1.75, 3.9, 3.05,
    "INFRASTRUCTURE (MONTHLY)",
    [
        ("VPS Server (4 vCPU, 8GB)", "$40/mo"),
        ("Domain & SSL Certificate", "$15/mo"),
        ("Email Service (SES/SendGrid)", "$20/mo"),
        ("Streaming CDN (basic tier)", "$50/mo"),
        ("AI API Costs (LLM tokens)", "$30/mo"),
        ("Backups & Monitoring", "$25/mo"),
    ],
    "Monthly Total", "$180/mo"
)

# Card 3: Third-Party
cost_card(s3, 8.9, 1.75, 3.9, 3.05,
    "THIRD-PARTY SERVICES",
    [
        ("Stripe Processing Fee", "2.9% + $0.30"),
        ("OAuth (Google/Apple)", "Free"),
        ("Push Notifications (Firebase)", "Free tier"),
        ("AI/LLM API (usage-based)", "~$0.01/query"),
        ("Streaming (usage-based)", "~$0.01/min"),
    ],
    "Variable", "Usage-based"
)

# Card 4: Total (highlighted)
cost_card(s3, 0.5, 5.0, 5.8, 1.55,
    "TOTAL INVESTMENT SUMMARY",
    [
        ("One-Time Development (8-10 wks)", "$8,500"),
        ("Infrastructure (3 months)", "$540"),
    ],
    "POC Total Investment", "~$10,000",
    highlight=True
)

# Card 5: Assumptions
card(s3, 6.6, 5.0, 6.2, 1.55)
rect(s3, 6.6, 5.0, 6.2, 0.4, ACCENT_LIGHT, CARD_BORDER)
txt(s3, 6.85, 5.07, 5, 0.28, "ASSUMPTIONS & NOTES", 11, ACCENT, bold=True)

assumptions = [
    "Development rate: ~$50/hr (mid-market full-stack rate)",
    "VPS: DigitalOcean/Hetzner tier with Docker pre-installed",
    "AI Guru Bot uses LLM API (e.g. Claude/OpenAI) with faith-specific prompts",
    "Excludes: Mobile apps, dedicated support, multi-region deployment",
    "Timeline assumes single developer, can accelerate with team",
]
ay = 5.48
for a in assumptions:
    # small dot
    dot = s3.shapes.add_shape(MSO_SHAPE.OVAL, Inches(6.9), Inches(ay + 0.05), Inches(0.05), Inches(0.05))
    dot.fill.solid()
    dot.fill.fore_color.rgb = ACCENT
    dot.line.fill.background()
    txt(s3, 7.05, ay, 5.5, 0.2, a, 9, SUBTLE)
    ay += 0.2

footer(s3, 3)


# ── Save ──
out = "/mnt/c/Projects/DAK/DAK_POC_Proposal_v2.pptx"
prs.save(out)
print(f"PPTX saved: {out}")
